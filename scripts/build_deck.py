#!/usr/bin/env python3
"""
Build the capstone viva deck from the course template.

    python3 scripts/build_deck.py

Reads  docs/capstone/presentation.pptx            (blank BITS template — never modified)
Writes docs/capstone/Claw_Capstone_Viva.pptx      (the deck to present)

Content comes from Capstone_Final_Submission.md per the slide-to-section map in §5.3.3.
Re-run it after the report changes.

Notes on the template, learned by inspecting it:
  * 4:3 (10 x 7.5 in). The body box on slides 2-10 is 9.29 in wide with only a
    0.27 in right margin, so long bullets wrap tightly. Keep them short.
  * Slide 1 has NO placeholders — six free-floating shapes at hardcoded positions.
    They are addressed by name and vertical position, and their runs are rewritten
    in place so the BITS brand purple (#9700FF / #9933FF) survives.
  * The template ships zero speaker notes. Every slide gets them here; the viva
    marks Q&A separately from the slides themselves.
"""

import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("python-pptx is not installed.  Run:  python3 -m pip install python-pptx lxml")

ROOT = Path(__file__).resolve().parent.parent
CAPSTONE = ROOT / "docs" / "capstone"
TEMPLATE = CAPSTONE / "presentation.pptx"
OUTPUT = CAPSTONE / "Claw_Capstone_Viva.pptx"
FIGURES = CAPSTONE / "figures"
SHOTS = CAPSTONE / "screenshots"
ASSETS = CAPSTONE / "deck-assets"         # generated illustrations
BUILD = CAPSTONE / ".deck-build"          # rasterised SVGs, gitignored

TITLE = {
    "project": "Claw — A Mobile-First Operations Hub\nfor Solo Service Providers",
    "team": "Om Mishra · Yash Agarwal · Divyanshu Pandey\n2023ebcs786 · 2023ebcs753 · 2023ebcs728",
    "course": "BSc Computer Science  ·  Group 82 (CodingNinjas)  ·  2025–26",
    "guide": "Swapnil Saurav",
}

# ---------------------------------------------------------------- tuning
# Everything a small layout tweak needs. Change a number here rather than
# editing the functions below.
LAYOUT = {
    "body_pt": 17,            # bullet size on a full-width slide
    "body_pt_narrow": 15,     # bullet size when an image sits beside the text
    "body_w_narrow": 5.15,    # inches — body box width when narrowed
    "img_right": (5.75, 2.05, 3.85, 4.55),   # left, top, width, max height (in)
    "img_below": (2.60, 4.30, 4.80, 2.90),
    "title_pt": 24,           # slide 1 project title
}

# Per-slide overrides, e.g. {9: {"body_pt": 15}} to shrink one dense slide.
TWEAKS: dict[int, dict] = {}

# (title, bullets, speaker notes, [(image_path, kind)])
#   kind "right"  -> image beside a narrowed body box
#   kind "below"  -> image under the bullets, full width
SLIDES = {
    2: (
        "Problem Statement",
        [
            "Solo providers run their business on four disconnected tools:",
            "WhatsApp · Calendar · Spreadsheet · memory",
            "Double bookings — negotiation is disconnected from the record",
            "Revenue leakage — unpaid sessions and lapsed clients slip through",
            "Client friction — 5–10 messages to agree one appointment",
            "Zero visibility — “what did I earn this month?” is unanswerable",
            "7.7 M gig workers in 2020–21 → 23.5 M projected by 2029–30",
        ],
        "Lead with the double-booking story — it is the most concrete failure and the one "
        "the whole product is built to prevent.\n\n"
        "The four failures are not independent: they all follow from the negotiation channel "
        "being disconnected from the record channel.\n\n"
        "Source for the workforce figures is NITI Aayog. If asked why this segment specifically: "
        "they cannot afford admin staff or per-seat SaaS, so operational friction converts "
        "directly into lost income.",
        [(ASSETS / "slide2-fragmented-tools.png", "right")],
    ),
    3: (
        "Objectives & Scope",
        [
            "Objectives",
            "Mobile-first CRM: client records with full appointment history",
            "Intelligent scheduling that removes multi-message negotiation",
            "Business analytics: revenue, trends, completion rate, inactive clients",
            "Recurring appointments — weekly, bi-weekly, monthly",
            "One-handed gesture interface for use between sessions",
            "Out of scope",
            "Payment processing · client self-booking portal · offline mode · app store release",
        ],
        "All five primary objectives were met, plus three secondary ones "
        "(WhatsApp reminders, inactive-client recovery, active/inactive service catalogue).\n\n"
        "The exclusions are deliberate product decisions, not things we ran out of time for. "
        "Client self-booking in particular: research showed this market still prefers to "
        "negotiate over WhatsApp, so a booking portal would have been built for a behaviour "
        "that does not exist yet.",
        [],
    ),
    4: (
        "Existing System / Literature Review",
        [
            "Generic schedulers — Calendly, Zoho Bookings",
            "Solve scheduling alone. No CRM, no client history, no revenue tracking",
            "Enterprise CRM — HubSpot, Pipedrive, Salesforce",
            "Built for sales teams; overwhelming for one practitioner managing appointments",
            "Vertical practice management — SimplePractice, TherapyNotes",
            "Well designed, but USD-priced and narrowed to healthcare compliance",
            "Gap: no integrated, affordable, mobile-first hub for the Indian solopreneur",
        ],
        "The point is not that these tools are bad — each is good at what it does. "
        "It is that a solo practitioner would need three of them, and the seams between "
        "them are exactly where the failures on the previous slide happen.\n\n"
        "If pushed on why not just use Calendly plus a spreadsheet: that IS the status quo, "
        "and it is what produces double bookings.",
        [],
    ),
    5: (
        "Proposed System Architecture",
        [
            "Three tiers: React Native client · Express REST API · PostgreSQL on Supabase",
            "Stateless JWT auth; token in Expo Secure Store (Keychain / Keystore)",
            "27 REST endpoints across six resources",
            "Every query filters by user_id — the whole authorisation model",
            "Six tables; composite index on (user_id, date_time) serves the three hottest reads",
        ],
        "Why a separate backend rather than calling Supabase directly with row-level security: "
        "recurring generation and conflict detection are far clearer in server-side JavaScript "
        "than in SQL policies, and the password-reset email needs a secret that cannot ship "
        "inside a mobile binary.\n\n"
        "Why store end_time rather than deriving it: it avoids a join on every read, and it "
        "keeps historical appointments correct after a service's duration or price is edited.",
        [(FIGURES / "fig-2.1-architecture.svg", "right")],
    ),
    6: (
        "Tools & Technologies",
        [
            "Mobile — React Native 0.81, Expo SDK 54, Expo Router 6, TypeScript",
            "UI — NativeWind 4, Reanimated, Gesture Handler, expo-haptics",
            "State — TanStack Query 5 (5-minute cache, optimistic updates)",
            "Backend — Node.js, Express 4, Zod validation, JWT, bcrypt",
            "Database — PostgreSQL on Supabase",
            "Quality — Jest (27 tests), ESLint, TypeScript strict, GitHub Actions CI",
        ],
        "Every choice has an alternative we rejected: Flutter (team knows TypeScript), "
        "MongoDB (the data is relational — clients own appointments own services), "
        "Firebase (no joins, weak aggregation, which would have made analytics painful).\n\n"
        "TanStack Query is the one worth calling out — its cache invalidation removed an "
        "entire class of stale-data bugs we were fighting by hand.",
        [],
    ),
    7: (
        "Implementation / Demo",
        [
            "Twelve modules, all fully implemented",
            "Booking reduced to four taps: client → service → date/time",
            "End time auto-calculated; overlap checked before every write",
            "Recurring series generated in one action, all-or-nothing on conflict",
            "Swipe right to complete, left to cancel — with haptic confirmation",
            "WhatsApp deep-links for reminders and win-back, no API subscription needed",
        ],
        "This is the slide to stop talking on and switch to the live device.\n\n"
        "Demo order that works: dashboard → book an appointment → deliberately book a "
        "clashing slot to show the conflict being refused → swipe one complete → analytics.\n\n"
        "The conflict refusal is the moment worth engineering the demo around; it is the "
        "product's central promise and takes ten seconds to show.",
        [(SHOTS / "03-dashboard.png", "right")],
    ),
    8: (
        "Results & Analysis",
        [
            "49 manual functional test cases — 100% pass",
            "27 automated tests over the scheduling core — 100% coverage, run on every push",
            "Dashboard loads ~1.2 s against a < 2 s target",
            "Search filters in ~200 ms against a < 500 ms target",
            "No crashes across a two-week test period; 100+ clients, 52-appointment series",
            "All six risks identified in Phase 1 successfully mitigated",
        ],
        "Distinguish the two kinds of testing if asked: 49 manual cases cover the system "
        "end to end on a real device; 27 automated tests cover the pure scheduling logic — "
        "the overlap rule and recurring date generation.\n\n"
        "Writing the automated tests found two real defects manual testing had missed: "
        "the 52-appointment cap refuses a weekly booking that runs a full calendar year "
        "(a year needs 53), and monthly series starting on the 31st overflow short months. "
        "Both are recorded in the report rather than quietly fixed.",
        [(SHOTS / "11-analytics.png", "right")],
    ),
    9: (
        "Challenges & Limitations",
        [
            "Recurring dates across month boundaries — 31 Jan has no equivalent in February",
            "Keeping the UI in step with the server after mutations; solved with query invalidation",
            "Scope discipline — MoSCoW prioritisation kept the core ahead of the extras",
            "No payment processing; practitioners still collect by UPI or cash",
            "No offline mode; the app needs connectivity",
            "Push notifications limited inside Expo Go — a platform constraint, not a defect",
        ],
        "Answer honestly on limitations. Each was a decision with a reason, and the report "
        "records them in §6.3.\n\n"
        "The strongest answer on payments: integrating a gateway is a compliance and "
        "reconciliation problem, not a coding one, and it would have consumed the time that "
        "went into making scheduling correct.\n\n"
        "If asked what we would do differently: write the automated tests first — they found "
        "bugs that fifty manual cases did not.",
        [],
    ),
    10: (
        "Conclusion & Future Work",
        [
            "A working end-to-end product, installable on a real device today",
            "All five primary and three secondary objectives delivered",
            "Short term — client self-booking page, UPI payments, production build and store release",
            "Medium term — offline-first sync, PDF invoices, Hindi and regional languages",
            "Long term — AI pricing and churn insight, team accounts, practitioner directory",
        ],
        "Close on the point that the most valuable feature is the least technically impressive: "
        "conflict detection is a single interval-overlap query, and it solves the most damaging "
        "problem these users have.\n\n"
        "The broader lesson was starting from the user's actual context rather than the "
        "technology — WhatsApp deep-links instead of building chat, INR by default, gestures "
        "designed for use between sessions.",
        [],
    ),
}


def rasterise(svg: Path, width_px: int = 1600) -> Path:
    """SVG -> PNG via rsvg-convert, so diagrams stay sharp when projected."""
    BUILD.mkdir(parents=True, exist_ok=True)
    png = BUILD / (svg.stem + ".png")
    if png.exists() and png.stat().st_mtime > svg.stat().st_mtime:
        return png
    try:
        subprocess.run(
            ["rsvg-convert", "-w", str(width_px), str(svg), "-o", str(png)],
            check=True, capture_output=True,
        )
        return png
    except (FileNotFoundError, subprocess.CalledProcessError):
        fallback = svg.with_suffix(".png")
        if fallback.exists():
            print(f"    rsvg-convert unavailable — using committed {fallback.name}")
            return fallback
        raise


def set_text(shape, text, size=None):
    """Rewrite a shape's text, preserving the first run's formatting."""
    tf = shape.text_frame
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    run = p0.runs[0] if p0.runs else p0.add_run()
    run.text = lines[0]
    if size:
        run.font.size = Pt(size)
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = run.font.size
        r.font.bold = run.font.bold
        r.font.name = run.font.name
        if run.font.color and run.font.color.type is not None:
            r.font.color.rgb = run.font.color.rgb


def fill_title_slide(slide):
    """Slide 1 has no placeholders; address the six shapes by name and position."""
    by_top = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: s.top,
    )
    # 1.94in "Presented by / Team Members", 3.89in title, 4.29in course, 5.27in guide
    for sh in by_top:
        top_in = sh.top / 914400
        if sh.name == "Title 1":
            set_text(sh, TITLE["project"], size=LAYOUT["title_pt"])
        elif top_in < 2.5:
            set_text(sh, "Presented by:\n" + TITLE["team"])
        elif 4.0 < top_in < 5.0:
            set_text(sh, TITLE["course"])
        elif top_in >= 5.0:
            set_text(sh, "Under the guidance of\n" + TITLE["guide"])


def fill_body(slide, bullets, narrowed=False, cfg=None):
    cfg = {**LAYOUT, **(cfg or {})}
    body = slide.placeholders[1]
    if narrowed:
        body.width = Inches(cfg["body_w_narrow"])
    tf = body.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    size = Pt(cfg["body_pt_narrow"] if narrowed else cfg["body_pt"])
    for i, text in enumerate(bullets):
        p = first if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = text
        run.font.size = size


def add_picture(slide, path, kind, cfg=None):
    cfg = {**LAYOUT, **(cfg or {})}
    l, t, w, h = cfg["img_right"] if kind == "right" else cfg["img_below"]
    pic = slide.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
    max_h = Inches(h)
    if pic.height > max_h:                       # re-fit by height, keep aspect
        ratio = max_h / pic.height
        pic.height, pic.width = int(pic.height * ratio), int(pic.width * ratio)
        pic.left = Inches(l) + int((Inches(w) - pic.width) / 2)   # re-centre
    return pic


def main():
    if not TEMPLATE.exists():
        sys.exit(f"Template not found: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))
    print(f"\n  template: {TEMPLATE.name}  ({len(prs.slides)} slides, "
          f"{prs.slide_width/914400:.0f}x{prs.slide_height/914400:.1f} in)\n")

    fill_title_slide(prs.slides[0])
    prs.slides[0].notes_slide.notes_text_frame.text = (
        "Introduce the team and the one-line pitch: Claw replaces the four disconnected tools "
        "a solo practitioner uses today with one app built for their phone.\n\n"
        "Keep this slide short — the marks are in the explanation, the demo and the questions."
    )
    print("   1  Title slide")

    for n, (title, bullets, notes, images) in SLIDES.items():
        slide = prs.slides[n - 1]
        slide.placeholders[0].text_frame.text = title
        cfg = TWEAKS.get(n)
        fill_body(slide, bullets, narrowed=bool(images), cfg=cfg)
        for path, kind in images:
            src = rasterise(path) if path.suffix == ".svg" else path
            if not src.exists():
                print(f"      missing image: {src}")
                continue
            add_picture(slide, src, kind, cfg=cfg)
        slide.notes_slide.notes_text_frame.text = notes
        img_note = f"  + {len(images)} image" if images else ""
        print(f"  {n:>2}  {title}{img_note}")

    prs.save(str(OUTPUT))
    size_kb = OUTPUT.stat().st_size / 1024
    print(f"\n  wrote {OUTPUT.relative_to(ROOT)}  ({size_kb:.0f} KB)")
    print(f"  template untouched: {TEMPLATE.relative_to(ROOT)}\n")
    print("  Open it and check every slide — there is no local pptx renderer,")
    print("  and the 4:3 body box wraps tightly.\n")


if __name__ == "__main__":
    main()
